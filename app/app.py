import streamlit as st
from pptx import Presentation
from pdfminer.high_level import extract_text

uploaded_file = st.file_uploader("Choose a Powerpoint or PDF file to upload")

if uploaded_file is not None:
    filename = uploaded_file.name

    if filename.endswith(".pptx"):
        prs = Presentation(uploaded_file)
        slides = prs.slides

        i=0
        # sort shapes on slide before reading
        for slide in slides:
            shapes = slide.shapes
            positions = []
            for shape in shapes:
                positions.append([shape.left, shape.top])

            shapes_sorted = (x for _,x in sorted(zip(positions, shapes), key = lambda x:(x[0][0], x[0][1])))

        # read shapes on slide
            st.write(f"Slide {i}:")
            i += 1
            for shape in shapes_sorted:
                if not shape.has_text_frame:
                    continue
                paragraphs = shape.text_frame.paragraphs
                for paragraph in paragraphs:
                    st.write(paragraph.text)
                st.write()

    if filename.endswith(".pdf"):
        text = extract_text(uploaded_file)
        st.write(text)
