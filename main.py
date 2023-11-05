from pptx import Presentation

prs = Presentation("data/test_slide_bcg.pptx")
slides = prs.slides

i=0
# sort shapes on slide before reading
for slide in slides:
    shapes = slide.shapes
    positions = []
    for shape in shapes:
        positions.append([shape.left, shape.top])

    shapes_sorted = (x for _,x in sorted(zip(positions, shapes), key = lambda x:(x[0][0], x[0][1])))

# read shapes on slide
    print(f"Slide {i}:")
    i += 1
    for shape in shapes_sorted:
        if not shape.has_text_frame:
            continue
        paragraphs = shape.text_frame.paragraphs
        for paragraph in paragraphs:
            print(paragraph.text)
        print()
